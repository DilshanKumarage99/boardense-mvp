import os
from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument

try:
    from openpyxl import load_workbook
except Exception:
    load_workbook = None

# Import Presentation from python-pptx optionally so the app can start
# even if python-pptx isn't installed. If missing, handle it when
# attempting to extract from .pptx files and provide a clear error.
try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import xlrd
except Exception:
    xlrd = None

def extract_text_from_document(file_path):
    """Extract text from various document formats"""
    
    ext = Path(file_path).suffix.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    elif ext in {'.xlsx', '.xls'}:
        return extract_text_from_excel(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def extract_text_from_excel(file_path):
    """Extract text from Excel workbooks (.xlsx/.xls)."""
    ext = Path(file_path).suffix.lower()

    if ext == '.xlsx':
        return extract_text_from_xlsx(file_path)
    if ext == '.xls':
        return extract_text_from_xls(file_path)

    raise ValueError(f"Unsupported Excel format: {ext}")


def extract_text_from_xlsx(file_path):
    """Extract text from XLSX using openpyxl."""
    if load_workbook is None:
        raise ImportError("openpyxl is not installed. Install it with: pip install openpyxl")

    output = []
    try:
        workbook_values = load_workbook(filename=file_path, data_only=True)
        workbook_formulas = load_workbook(filename=file_path, data_only=False)

        for sheet_values, sheet_formulas in zip(workbook_values.worksheets, workbook_formulas.worksheets):
            sheet_lines = [f"\n--- SHEET: {sheet_values.title} ---"]

            max_row = max(sheet_values.max_row, sheet_formulas.max_row)
            max_col = max(sheet_values.max_column, sheet_formulas.max_column)

            for row_index in range(1, max_row + 1):
                row_values = []
                for col_index in range(1, max_col + 1):
                    value_cell = sheet_values.cell(row=row_index, column=col_index).value
                    formula_cell = sheet_formulas.cell(row=row_index, column=col_index).value

                    cell_value = value_cell if value_cell is not None else formula_cell
                    if cell_value is None:
                        continue

                    cell_text = str(cell_value).strip()
                    if cell_text:
                        row_values.append(cell_text)

                if row_values:
                    sheet_lines.append(" | ".join(row_values))

            if len(sheet_lines) > 1:
                output.extend(sheet_lines)
    except Exception as e:
        raise ValueError(f"Error reading XLSX: {e}")

    return "\n".join(output)


def extract_text_from_xls(file_path):
    """Extract text from XLS using xlrd."""
    if xlrd is None:
        raise ImportError("xlrd is not installed. Install it with: pip install xlrd")

    output = []
    try:
        workbook = xlrd.open_workbook(file_path)
        for sheet in workbook.sheets():
            sheet_lines = [f"\n--- SHEET: {sheet.name} ---"]
            for row_index in range(sheet.nrows):
                row = sheet.row_values(row_index)
                row_values = [str(cell).strip() for cell in row if str(cell).strip()]
                if row_values:
                    sheet_lines.append(" | ".join(row_values))
            if len(sheet_lines) > 1:
                output.extend(sheet_lines)
    except Exception as e:
        raise ValueError(f"Error reading XLS: {e}")

    return "\n".join(output)

def extract_text_from_pdf(file_path):
    """Extract text from PDF"""
    text = []
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
    except Exception as e:
        raise ValueError(f"Error reading PDF: {e}")
    
    return '\n'.join(text)

def extract_text_from_docx(file_path):
    """Extract text from DOCX"""
    text = []
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
    except Exception as e:
        raise ValueError(f"Error reading DOCX: {e}")
    
    return '\n'.join(text)

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    if Presentation is None:
        raise ImportError(
            "python-pptx is not installed. Install it with: pip install python-pptx"
        )
    text = []
    try:
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"\n--- SLIDE {slide_num} ---"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_data.append(cell.text)
                        if row_data:
                            slide_text.append(" | ".join(row_data))
            
            # Add slide text if found
            if len(slide_text) > 1:
                text.extend(slide_text)
    
    except Exception as e:
        raise ValueError(f"Error reading PPTX: {e}")
    
    return '\n'.join(text) if text else ""

def process_document(document):
    """Process uploaded document and extract text"""
    try:
        text = extract_text_from_document(document.file_path)
        document.content_extracted = text
        from app import db
        db.session.commit()
        return True
    except Exception as e:
        print(f"Error processing document {document.id}: {e}")
        return False
